import subprocess
from pathlib import Path
from PIL import Image

from backend.downloader import ffmpeg_path

IMAGE_FORMATS = {"jpg", "jpeg", "png", "webp", "bmp", "gif", "tiff", "ico"}
AUDIO_FORMATS = {"mp3", "flac", "wav", "m4a", "ogg", "opus", "aac"}
VIDEO_FORMATS = {"mp4", "avi", "mov", "mkv", "webm", "gif"}
DOC_FORMATS = {"pdf", "pptx"}

ALL_TARGETS = IMAGE_FORMATS | AUDIO_FORMATS | VIDEO_FORMATS | DOC_FORMATS

ICO_SIZES = [16, 24, 32, 48, 64, 128, 256]


def detect_type(path: Path) -> str:
    ext = path.suffix.lstrip(".").lower()
    if ext in IMAGE_FORMATS:
        return "image"
    if ext in AUDIO_FORMATS:
        return "audio"
    if ext in VIDEO_FORMATS or ext in {"flv", "wmv", "m4v", "3gp"}:
        return "video"
    if ext == "pdf":
        return "pdf"
    if ext == "pptx":
        return "pptx"
    return "unknown"


def get_targets(kind: str) -> list:
    if kind == "image":
        return sorted(IMAGE_FORMATS | {"pdf", "pptx"})
    if kind == "audio":
        return sorted(AUDIO_FORMATS)
    if kind == "video":
        return sorted(VIDEO_FORMATS | AUDIO_FORMATS)
    if kind == "pdf":
        return sorted(IMAGE_FORMATS | {"pptx"})
    if kind == "pptx":
        return ["pdf"] + sorted(IMAGE_FORMATS - {"ico"})
    return []


def _convert_image(src: Path, dst: Path, target: str):
    img = Image.open(src)

    if target == "ico":
        rgba = img.convert("RGBA")
        sizes = [(s, s) for s in ICO_SIZES if s <= max(rgba.size)] or [(32, 32)]
        rgba.save(dst, format="ICO", sizes=sizes)
        return

    if target in ("jpg", "jpeg"):
        img = img.convert("RGB")

    fmt = "JPEG" if target in ("jpg", "jpeg") else target.upper()
    img.save(dst, format=fmt)


def _convert_image_to_pdf(src: Path, dst: Path):
    import img2pdf
    with open(dst, "wb") as f:
        f.write(img2pdf.convert(str(src)))


EMU_PER_PX = 9525
EMU_MIN = 914400       # 1 inch
EMU_MAX = 51206400     # 56 inches


def _slide_dims(w_px: int, h_px: int) -> tuple[int, int]:
    w = w_px * EMU_PER_PX
    h = h_px * EMU_PER_PX
    smallest = min(w, h)
    largest = max(w, h)
    if smallest < EMU_MIN:
        scale = EMU_MIN / smallest
        w, h = int(w * scale), int(h * scale)
    elif largest > EMU_MAX:
        scale = EMU_MAX / largest
        w, h = int(w * scale), int(h * scale)
    return w, h


def _convert_image_to_pptx(src: Path, dst: Path):
    from pptx import Presentation
    from pptx.util import Emu

    img = Image.open(src)
    w_px, h_px = img.size
    prs = Presentation()
    prs.slide_width, prs.slide_height = _slide_dims(w_px, h_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(src), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(dst)


def _convert_pdf_to_images_or_pptx(src: Path, dst: Path, target: str):
    import fitz  # pymupdf

    doc = fitz.open(src)

    if target == "pptx":
        from pptx import Presentation
        import tempfile

        prs = Presentation()
        with tempfile.TemporaryDirectory() as tmp:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                page_img = Path(tmp) / f"page_{i}.png"
                pix.save(page_img)
                if i == 0:
                    prs.slide_width, prs.slide_height = _slide_dims(pix.width, pix.height)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(str(page_img), 0, 0, width=prs.slide_width, height=prs.slide_height)
            prs.save(dst)
        return

    # single page -> single image; only first page is used for image targets
    page = doc[0]
    pix = page.get_pixmap(dpi=200)
    pix.save(dst)


def _convert_pptx_to_pdf(src: Path, dst: Path):
    raise RuntimeError(
        "PPTX -> PDF/картинка требует LibreOffice или PowerPoint (рендер слайдов) — "
        "не поддерживается без тяжёлых внешних зависимостей."
    )


def _convert_audio_video(src: Path, dst: Path, target: str):
    ffmpeg_dir = ffmpeg_path()
    exe = str(Path(ffmpeg_dir) / "ffmpeg.exe") if ffmpeg_dir else "ffmpeg"

    args = [exe, "-y", "-i", str(src)]

    if target == "gif":
        args += ["-vf", "fps=12,scale=480:-1:flags=lanczos"]
    elif target in AUDIO_FORMATS:
        args += ["-vn"]
        if target == "mp3":
            args += ["-codec:a", "libmp3lame", "-b:a", "320k"]
        elif target == "flac":
            args += ["-codec:a", "flac"]
        elif target == "wav":
            args += ["-codec:a", "pcm_s16le"]
        elif target == "aac":
            args += ["-codec:a", "aac", "-b:a", "256k"]

    args.append(str(dst))

    result = subprocess.run(args, capture_output=True, text=True, timeout=600)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg error: {result.stderr[-800:]}")


def convert(src: Path, dst: Path, target: str) -> tuple[bool, str | None]:
    try:
        kind = detect_type(src)

        if kind == "image":
            if target == "pdf":
                _convert_image_to_pdf(src, dst)
            elif target == "pptx":
                _convert_image_to_pptx(src, dst)
            else:
                _convert_image(src, dst, target)

        elif kind == "pdf":
            _convert_pdf_to_images_or_pptx(src, dst, target)

        elif kind == "pptx":
            _convert_pptx_to_pdf(src, dst)

        elif kind in ("audio", "video"):
            _convert_audio_video(src, dst, target)

        else:
            return False, f"неизвестный формат: {src.suffix}"

        return True, None
    except Exception as e:
        return False, str(e)
